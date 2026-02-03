import os
import logging
from typing import List, Optional
from tqdm import tqdm

# --- 格式处理库 ---
import pymupdf4llm  # PDF 神器
from docx import Document
from pptx import Presentation
import pandas as pd

# --- LangChain 组件 ---
from langchain_core.documents import Document as LangchainDocument
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma

import torch

# --- 配置 ---
SOURCE_DIR = r"./attachments"  # 你的课件存放目录
DB_DIR = "./chroma_db"  # 向量数据库路径
CHUNK_SIZE = 800  # 分块大小
CHUNK_OVERLAP = 100  # 重叠部分

# --- 日志配置 ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


class DocumentConverter:
    """
    多格式转换器：将二进制文件统一转换为清洗后的 Markdown 文本
    """

    @staticmethod
    def convert_pdf(file_path: str) -> str:
        """使用 PyMuPDF4LLM 将 PDF 转换为 Markdown (保留表格结构)"""
        try:
            # pymupdf4llm 直接返回 markdown 字符串
            md_text = pymupdf4llm.to_markdown(file_path)
            return md_text
        except Exception as e:
            logging.error(f"❌ PDF Convert Error ({file_path}): {e}")
            return ""

    @staticmethod
    def convert_docx(file_path: str) -> str:
        """提取 Word 文档并保留基本结构"""
        try:
            doc = Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    # 简单的标题识别逻辑
                    if para.style.name.startswith('Heading'):
                        full_text.append(f"## {para.text}")
                    else:
                        full_text.append(para.text)
            return "\n\n".join(full_text)
        except Exception as e:
            logging.error(f"❌ DOCX Convert Error ({file_path}): {e}")
            return ""

    @staticmethod
    def convert_pptx(file_path: str) -> str:
        """提取 PPT 内容，按幻灯片分页"""
        try:
            prs = Presentation(file_path)
            full_text = []
            for i, slide in enumerate(prs.slides):
                slide_content = [f"## Slide {i + 1}"]

                # 尝试提取标题
                if slide.shapes.title:
                    slide_content.append(f"### {slide.shapes.title.text}")

                # 提取正文文本框
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # 避免重复标题
                        if shape == slide.shapes.title:
                            continue
                        slide_content.append(shape.text)

                full_text.append("\n".join(slide_content))
            return "\n\n---\n\n".join(full_text)
        except Exception as e:
            logging.error(f"❌ PPTX Convert Error ({file_path}): {e}")
            return ""

    @staticmethod
    def convert_excel(file_path: str) -> str:
        """
        通用 Excel 转换器：支持 .xlsx (标准), .xlsm (带宏), .xlsb (二进制)
        """
        try:
            ext = os.path.splitext(file_path)[1].lower()

            # 1. 智能选择引擎
            engine = None
            if ext == '.xlsb':
                engine = 'pyxlsb'  # 二进制专用引擎
            else:
                engine = 'openpyxl'  # .xlsx 和 .xlsm 用这个

            # 2. 加载文件
            xls = pd.ExcelFile(file_path, engine=engine)
            full_text = []

            for sheet_name in xls.sheet_names:
                # 读取数据 (自动忽略 .xlsm 中的 VBA 代码)
                df = pd.read_excel(xls, sheet_name=sheet_name)

                # 3. 数据清洗 (这是我们之前优化的核心)
                df = df.fillna("")  # 清洗 NaN

                # 截断过大的表格 (防止 Token 爆炸)
                if len(df) > 50:
                    df = df.head(50)
                    full_text.append(f"> [!WARNING] Table truncated (showing first 50 rows)")

                if not df.empty:
                    md_table = df.to_markdown(index=False)
                    full_text.append(f"## Sheet: {sheet_name}\n\n{md_table}")

            return "\n\n".join(full_text)

        except ImportError as e:
            if '.xlsb' in file_path:
                logging.error(f"❌ Missing Library: Please run `uv pip install pyxlsb` to read .xlsb files.")
            return ""
        except Exception as e:
            logging.error(f"❌ Excel Convert Error ({file_path}): {e}")
            return ""


class KnowledgeIndexer:
    def __init__(self):
        if torch.cuda.is_available():
            device_type = 'cuda'
        elif torch.backends.mps.is_available():
            device_type = 'mps'  # Apple Silicon 的加速器
        else:
            device_type = 'cpu'

        self.embeddings = HuggingFaceEmbeddings(
            model_name="BAAI/bge-m3",
            model_kwargs={'device': device_type}
        )

        self.converter = DocumentConverter()

        # 文本分块器 (针对 Markdown 优化)
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            separators=["\n## ", "\n### ", "\n", " ", ""]  # 优先按标题切分
        )

    def process_directory(self, source_dir: str):
        """扫描并处理目录下所有支持的文件"""
        if not os.path.exists(source_dir):
            os.makedirs(source_dir)
            print(f"📂 Created directory: {source_dir}. Put your files here!")
            return

        all_docs = []
        files = [f for f in os.listdir(source_dir) if not f.startswith("~")]  # 忽略临时文件

        print(f"🔍 Found {len(files)} files. Starting conversion...")

        for filename in tqdm(files, desc="Converting"):
            file_path = os.path.join(source_dir, filename)
            ext = os.path.splitext(filename)[1].lower()

            content = ""
            if ext == ".pdf":
                content = self.converter.convert_pdf(file_path)
            elif ext == ".docx":
                content = self.converter.convert_docx(file_path)
            elif ext == ".pptx":
                content = self.converter.convert_pptx(file_path)
            elif ext in [".xlsx", ".xlsm", ".xlsb", ".xls"]:
                content = self.converter.convert_excel(file_path)
            elif ext == ".md":
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
            else:
                logging.warning(f"⚠️ Skipped unsupported format: {filename}")
                continue

            if content:
                # 封装为 LangChain Document，带上元数据
                doc = LangchainDocument(
                    page_content=content,
                    metadata={"source": filename, "type": ext}
                )
                all_docs.append(doc)

        if not all_docs:
            print("⚠️ No valid content extracted.")
            return

        # 分块
        print(f"✂️ Splitting {len(all_docs)} documents...")
        chunks = self.splitter.split_documents(all_docs)
        print(f"🧩 Generated {len(chunks)} chunks.")

        # 存入数据库
        print("💾 Persisting to ChromaDB (this may take a while)...")
        vectordb = Chroma.from_documents(
            documents=chunks,
            embedding=self.embeddings,
            persist_directory=DB_DIR,
            collection_name="fintech_knowledge"
        )
        print("✅ Indexing Complete! Your agent can now read your course materials.")


if __name__ == "__main__":
    indexer = KnowledgeIndexer()
    indexer.process_directory(SOURCE_DIR)